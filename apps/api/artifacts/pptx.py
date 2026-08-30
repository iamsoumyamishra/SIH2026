"""PPTX artifact generation using python-pptx."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def generate_pptx(
    path: str | Path,
    title: str,
    slides: list[dict],
) -> Path:
    """slides: list of {"title": str, "bullets": [str]}."""
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    layout = prs.slide_layouts[1]

    first = prs.slides.add_slide(prs.slide_layouts[0])
    first.shapes.title.text = title

    for slide in slides:
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide.get("title", "")
        body = s.placeholders[1]
        tf = body.text_frame
        tf.text = ""
        for i, bullet in enumerate(slide.get("bullets", [])):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet

    prs.save(str(path))
    return path
